"""Document text extractor for external files (Word, PDF, Markdown, etc.).

Extracts plain text from various document formats to use as context
for AI paper generation.
"""

from __future__ import annotations

from pathlib import Path


def extract_text(file_path: str | Path, max_chars: int = 20000) -> str:
    """Extract text from a document file.

    Supports: .docx, .pdf, .md, .txt, .rst, .tex, .log, .pptx
    Returns extracted text (truncated to max_chars).
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()

    if ext == ".docx":
        return _extract_docx(path, max_chars)
    elif ext == ".pdf":
        return _extract_pdf(path, max_chars)
    elif ext == ".pptx":
        return _extract_pptx(path, max_chars)
    else:
        return _extract_plaintext(path, max_chars)


def _extract_plaintext(path: Path, max_chars: int) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")[:max_chars]
    except Exception as e:
        return f"[Error reading {path.name}: {e}]"


def _extract_docx(path: Path, max_chars: int) -> str:
    try:
        from docx import Document
    except ImportError:
        return (
            f"[Cannot read {path.name}: python-docx not installed. "
            "Run: pip install python-docx]"
        )

    try:
        doc = Document(str(path))
        paragraphs: list[str] = []
        total = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
                total += len(text)
                if total >= max_chars:
                    break

        # Also extract tables
        for table in doc.tables:
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                paragraphs.append("Table:\n" + "\n".join(rows))
                total += sum(len(r) for r in rows)
                if total >= max_chars:
                    break

        return "\n\n".join(paragraphs)[:max_chars]
    except Exception as e:
        return f"[Error reading {path.name}: {e}]"


def _extract_pdf(path: Path, max_chars: int) -> str:
    # Try PyPDF2 first, then pdfplumber
    text = _try_pypdf2(path, max_chars)
    if text and len(text.strip()) > 50:
        return text

    text = _try_pdfplumber(path, max_chars)
    if text:
        return text

    return f"[Cannot extract text from {path.name}: install PyPDF2 or pdfplumber]"


def _try_pypdf2(path: Path, max_chars: int) -> str:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return ""

    try:
        reader = PdfReader(str(path))
        pages: list[str] = []
        total = 0
        for page in reader.pages:
            text = page.extract_text() or ""
            pages.append(text)
            total += len(text)
            if total >= max_chars:
                break
        return "\n\n".join(pages)[:max_chars]
    except Exception:
        return ""


def _try_pdfplumber(path: Path, max_chars: int) -> str:
    try:
        import pdfplumber
    except ImportError:
        return ""

    try:
        pages: list[str] = []
        total = 0
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                pages.append(text)
                total += len(text)
                if total >= max_chars:
                    break
        return "\n\n".join(pages)[:max_chars]
    except Exception:
        return ""


def _extract_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return (
            f"[Cannot read {path.name}: python-pptx not installed. "
            "Run: pip install python-pptx]"
        )

    try:
        prs = Presentation(str(path))
        slides: list[str] = []
        total = 0
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
                total += sum(len(t) for t in texts)
                if total >= max_chars:
                    break
        return "\n\n".join(slides)[:max_chars]
    except Exception as e:
        return f"[Error reading {path.name}: {e}]"
